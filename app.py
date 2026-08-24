"""
Study Buddy AI Agent
-------------------------------------------------------
Main Streamlit Application

Supports:
    PDF
    TXT
    CSV
    PPT
    PPTX

Features:
    - File upload
    - Text extraction
    - Summarization
    - Important topics
    - Quiz generation
    - Quiz evaluation
    - Flashcards
    - Revision planner
    - Topic explanation
    - Progress tracking
    - Analytics charts

Compatible with:
    Python 3.14.7
"""

import os
import io
import json
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt

from PyPDF2 import PdfReader
from pptx import Presentation


# =========================================================
# PATHS
# =========================================================

UPLOAD_FOLDER = "uploads"
DATA_FOLDER = "data"
TEMP_FOLDER = "temp"

PROGRESS_FILE = os.path.join(
    DATA_FOLDER,
    "progress.json"
)

QUIZ_HISTORY_FILE = os.path.join(
    DATA_FOLDER,
    "quiz_history.json"
)

SCHEDULE_FILE = os.path.join(
    DATA_FOLDER,
    "study_schedule.json"
)

CURRENT_QUIZ_FILE = os.path.join(
    DATA_FOLDER,
    "current_quiz.json"
)

EXTRACTED_TEXT_FILE = os.path.join(
    TEMP_FOLDER,
    "extracted_text.txt"
)


# =========================================================
# CREATE REQUIRED FOLDERS
# =========================================================

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(DATA_FOLDER, exist_ok=True)
os.makedirs(TEMP_FOLDER, exist_ok=True)


# =========================================================
# DEFAULT JSON FILES
# =========================================================

DEFAULT_PROGRESS = {
    "study_days": 0,
    "total_files_uploaded": 0,
    "completed_topics": [],
    "quiz_scores": [],
    "average_score": 0,
    "best_score": 0,
    "last_score": 0,
    "total_quizzes": 0,
    "study_time_minutes": 0,
    "files": [],
    "last_study_date": ""
}


def create_json_file(path, default_data):
    """Create JSON file if it does not exist."""

    if not os.path.exists(path):

        with open(
            path,
            "w",
            encoding="utf-8"
        ) as file:

            json.dump(
                default_data,
                file,
                indent=4,
                ensure_ascii=False
            )


create_json_file(
    PROGRESS_FILE,
    DEFAULT_PROGRESS
)

create_json_file(
    QUIZ_HISTORY_FILE,
    []
)

create_json_file(
    SCHEDULE_FILE,
    []
)

create_json_file(
    CURRENT_QUIZ_FILE,
    []
)


# =========================================================
# PAGE CONFIGURATION
# =========================================================

st.set_page_config(
    page_title="Study Buddy AI Agent",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)


# =========================================================
# CUSTOM CSS
# =========================================================

st.markdown(
    """
    <style>

    .main-title {
        font-size: 42px;
        font-weight: 700;
        margin-bottom: 5px;
    }

    .subtitle {
        font-size: 18px;
        margin-bottom: 25px;
    }

    .metric-card {
        padding: 15px;
        border-radius: 10px;
        border: 1px solid #ddd;
        text-align: center;
    }

    </style>
    """,
    unsafe_allow_html=True
)


# =========================================================
# SESSION STATE
# =========================================================

if "extracted_text" not in st.session_state:

    st.session_state.extracted_text = ""

if "uploaded_files" not in st.session_state:

    st.session_state.uploaded_files = []

if "current_file" not in st.session_state:

    st.session_state.current_file = ""

if "quiz_generated" not in st.session_state:

    st.session_state.quiz_generated = False

if "quiz_result" not in st.session_state:

    st.session_state.quiz_result = None


# =========================================================
# JSON HELPERS
# =========================================================

def load_json(path, default):

    try:

        if not os.path.exists(path):
            return default

        with open(
            path,
            "r",
            encoding="utf-8"
        ) as file:

            return json.load(file)

    except Exception:

        return default


def save_json(path, data):

    with open(
        path,
        "w",
        encoding="utf-8"
    ) as file:

        json.dump(
            data,
            file,
            indent=4,
            ensure_ascii=False
        )


# =========================================================
# PDF READER
# =========================================================

def read_pdf(file_bytes):

    text_parts = []

    pdf = PdfReader(
        io.BytesIO(file_bytes)
    )

    for page in pdf.pages:

        text = page.extract_text()

        if text:

            text_parts.append(text)

    return "\n".join(text_parts)


# =========================================================
# TEXT READER
# =========================================================

def read_txt(file_bytes):

    return file_bytes.decode(
        "utf-8",
        errors="ignore"
    )


# =========================================================
# CSV READER
# =========================================================

def read_csv(file_bytes):

    dataframe = pd.read_csv(
        io.BytesIO(file_bytes)
    )

    return dataframe.to_string(
        index=False
    )


# =========================================================
# POWERPOINT READER
# =========================================================

def read_ppt(file_bytes):

    presentation = Presentation(
        io.BytesIO(file_bytes)
    )

    text_parts = []

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        text_parts.append(
            f"\nSlide {slide_number}\n"
        )

        for shape in slide.shapes:

            if hasattr(
                shape,
                "text"
            ):

                if shape.text.strip():

                    text_parts.append(
                        shape.text
                    )

    return "\n".join(text_parts)


# =========================================================
# UNIVERSAL FILE READER
# =========================================================

def read_uploaded_file(uploaded_file):

    filename = uploaded_file.name

    extension = os.path.splitext(
        filename
    )[1].lower()

    file_bytes = uploaded_file.getvalue()

    if extension == ".pdf":

        return read_pdf(file_bytes)

    elif extension == ".txt":

        return read_txt(file_bytes)

    elif extension == ".csv":

        return read_csv(file_bytes)

    elif extension in [".ppt", ".pptx"]:

        return read_ppt(file_bytes)

    else:

        raise ValueError(
            "Unsupported file type."
        )


# =========================================================
# SAVE UPLOADED FILE
# =========================================================

def save_uploaded_file(uploaded_file):

    file_path = os.path.join(
        UPLOAD_FOLDER,
        uploaded_file.name
    )

    with open(
        file_path,
        "wb"
    ) as file:

        file.write(
            uploaded_file.getbuffer()
        )

    return file_path


# =========================================================
# SAVE EXTRACTED TEXT
# =========================================================

def save_extracted_text(text):

    with open(
        EXTRACTED_TEXT_FILE,
        "w",
        encoding="utf-8"
    ) as file:

        file.write(text)


# =========================================================
# PROGRESS FUNCTIONS
# =========================================================

def add_uploaded_file_to_progress(
    filename
):

    progress = load_json(
        PROGRESS_FILE,
        DEFAULT_PROGRESS.copy()
    )

    if filename not in progress["files"]:

        progress["files"].append(
            filename
        )

    progress[
        "total_files_uploaded"
    ] = len(
        progress["files"]
    )

    save_json(
        PROGRESS_FILE,
        progress
    )


def update_quiz_progress(score):

    progress = load_json(
        PROGRESS_FILE,
        DEFAULT_PROGRESS.copy()
    )

    if "quiz_scores" not in progress:
        progress["quiz_scores"] = []

    progress[
        "quiz_scores"
    ].append(score)

    progress[
        "last_score"
    ] = score

    progress[
        "total_quizzes"
    ] = len(
        progress["quiz_scores"]
    )

    progress[
        "average_score"
    ] = round(
        sum(progress["quiz_scores"])
        /
        len(progress["quiz_scores"]),
        2
    )

    progress[
        "best_score"
    ] = max(
        progress["quiz_scores"]
    )

    save_json(
        PROGRESS_FILE,
        progress
    )


# =========================================================
# IMPORT PROJECT MODULES
# =========================================================

def import_project_modules():

    modules = {}

    try:
        from summarizer import summarize
        modules["summarize"] = summarize
    except ImportError:
        pass

    try:
        from topics import extract_topics
        modules["extract_topics"] = extract_topics
    except ImportError:
        pass

    try:
        from quiz import generate, load_quiz
        modules["quiz_generate"] = generate
        modules["load_quiz"] = load_quiz
    except ImportError:
        pass

    try:
        from evaluator import evaluate
        modules["evaluate"] = evaluate
    except ImportError:
        pass

    try:
        from flashcards import create_flashcards
        modules["flashcards"] = create_flashcards
    except ImportError:
        pass

    try:
        from planner import create_plan
        modules["planner"] = create_plan
    except ImportError:
        pass

    try:
        from explain import explain_topic
        modules["explain"] = explain_topic
    except ImportError:
        pass

    return modules


modules = import_project_modules()


# =========================================================
# SIDEBAR
# =========================================================

with st.sidebar:

    st.markdown(
        "# 📚 Study Buddy"
    )

    st.caption(
        "AI-powered study assistant"
    )

    st.divider()

    page = st.radio(

        "Navigation",

        [
            "🏠 Home",
            "📤 Upload Materials",
            "📝 Summary",
            "🔑 Important Topics",
            "🧠 Quiz",
            "🗂️ Flashcards",
            "📅 Revision Planner",
            "💡 Explain Topic",
            "📊 Progress & Analytics"
        ]

    )

    st.divider()

    st.caption(
        "Python 3.14.7"
    )


# =========================================================
# HEADER
# =========================================================

st.markdown(
    '<div class="main-title">📚 Study Buddy AI Agent</div>',
    unsafe_allow_html=True
)

st.markdown(
    '<div class="subtitle">'
    "Upload your study materials and turn them into "
    "summaries, topics, quizzes, flashcards and revision plans."
    "</div>",
    unsafe_allow_html=True
)


# =========================================================
# HOME
# =========================================================

if page == "🏠 Home":

    progress = load_json(
        PROGRESS_FILE,
        DEFAULT_PROGRESS
    )

    col1, col2, col3, col4 = st.columns(4)

    with col1:

        st.metric(
            "📚 Files",
            progress.get(
                "total_files_uploaded",
                0
            )
        )

    with col2:

        st.metric(
            "📝 Quizzes",
            progress.get(
                "total_quizzes",
                0
            )
        )

    with col3:

        st.metric(
            "🏆 Best Score",
            f"{progress.get('best_score', 0)}%"
        )

    with col4:

        st.metric(
            "📈 Average",
            f"{progress.get('average_score', 0)}%"
        )

    st.divider()

    st.subheader(
        "What can Study Buddy do?"
    )

    features = [

        "📤 Upload PDF, TXT, CSV, PPT or PPTX",

        "📝 Generate study summaries",

        "🔑 Extract important topics",

        "🧠 Generate quizzes",

        "🗂️ Create flashcards",

        "📅 Create revision plans",

        "💡 Explain concepts",

        "📊 Track your study progress"

    ]

    for feature in features:

        st.write(
            f"✅ {feature}"
        )

    if not st.session_state.extracted_text:

        st.info(
            "Start by going to "
            "**Upload Materials**."
        )


# =========================================================
# UPLOAD MATERIALS
# =========================================================

elif page == "📤 Upload Materials":

    st.subheader(
        "📤 Upload Study Materials"
    )

    st.write(
        "Supported formats: "
        "**PDF, TXT, CSV, PPT, PPTX**"
    )

    uploaded_files = st.file_uploader(

        "Choose your study files",

        type=[
            "pdf",
            "txt",
            "csv",
            "ppt",
            "pptx"
        ],

        accept_multiple_files=True

    )

    if uploaded_files:

        combined_text = []

        for uploaded_file in uploaded_files:

            try:

                save_uploaded_file(
                    uploaded_file
                )

                extracted = read_uploaded_file(
                    uploaded_file
                )

                if extracted.strip():

                    combined_text.append(

                        f"\n\n===== "
                        f"{uploaded_file.name} "
                        f"=====\n\n"
                        f"{extracted}"

                    )

                    add_uploaded_file_to_progress(
                        uploaded_file.name
                    )

                else:

                    st.warning(
                        f"No text could be extracted "
                        f"from {uploaded_file.name}"
                    )

            except Exception as error:

                st.error(
                    f"Could not read "
                    f"{uploaded_file.name}: "
                    f"{error}"
                )

        final_text = "\n".join(
            combined_text
        )

        if final_text.strip():

            st.session_state.extracted_text = (
                final_text
            )

            st.session_state.uploaded_files = [
                file.name
                for file in uploaded_files
            ]

            st.session_state.current_file = (
                uploaded_files[0].name
            )

            save_extracted_text(
                final_text
            )

            st.success(
                f"Successfully processed "
                f"{len(uploaded_files)} file(s)."
            )

            st.info(
                f"Extracted approximately "
                f"{len(final_text.split())} words."
            )

    if st.session_state.extracted_text:

        st.divider()

        st.subheader(
            "Uploaded Files"
        )

        for filename in st.session_state.uploaded_files:

            st.write(
                f"📄 {filename}"
            )

        with st.expander(
            "Preview Extracted Text"
        ):

            st.text_area(
                "Extracted content",
                st.session_state.extracted_text,
                height=400
            )

        st.download_button(

            "⬇️ Download Extracted Text",

            data=st.session_state.extracted_text,

            file_name="extracted_study_material.txt",

            mime="text/plain"

        )


# =========================================================
# SUMMARY
# =========================================================

elif page == "📝 Summary":

    st.subheader(
        "📝 Study Summary"
    )

    text = st.session_state.extracted_text

    if not text:

        st.warning(
            "Please upload study material first."
        )

    else:

        if "summarize" in modules:

            if st.button(
                "Generate Summary"
            ):

                try:

                    summary = modules[
                        "summarize"
                    ](text)

                    st.success(
                        "Summary generated."
                    )

                    st.write(
                        summary
                    )

                except Exception as error:

                    st.error(
                        f"Summary error: {error}"
                    )

        else:

            st.error(
                "summarizer.py could not be loaded."
            )


# =========================================================
# IMPORTANT TOPICS
# =========================================================

elif page == "🔑 Important Topics":

    st.subheader(
        "🔑 Important Study Topics"
    )

    text = st.session_state.extracted_text

    if not text:

        st.warning(
            "Please upload study material first."
        )

    elif "extract_topics" not in modules:

        st.error(
            "topics.py could not be loaded."
        )

    else:

        if st.button(
            "Extract Important Topics"
        ):

            try:

                topics = modules[
                    "extract_topics"
                ](text)

                st.success(
                    "Important topics extracted."
                )

                if isinstance(
                    topics,
                    list
                ):

                    for number, topic in enumerate(
                        topics,
                        start=1
                    ):

                        st.write(
                            f"**{number}.** {topic}"
                        )

                else:

                    st.write(topics)

            except Exception as error:

                st.error(
                    f"Topic extraction error: {error}"
                )


# =========================================================
# QUIZ
# =========================================================

elif page == "🧠 Quiz":

    st.subheader(
        "🧠 Study Quiz"
    )

    text = st.session_state.extracted_text

    if not text:

        st.warning(
            "Please upload study material first."
        )

    elif "quiz_generate" not in modules:

        st.error(
            "quiz.py could not be loaded."
        )

    else:

        number_of_questions = st.slider(

            "Number of questions",

            min_value=5,

            max_value=30,

            value=20

        )

        if st.button(
            "🆕 Generate New Quiz"
        ):

            try:

                modules[
                    "quiz_generate"
                ](
                    text,
                    number_of_questions
                )

                st.session_state.quiz_generated = True
                st.session_state.quiz_result = None

                st.success(
                    "New quiz generated successfully!"
                )

                st.rerun()

            except Exception as error:

                st.error(
                    f"Quiz generation error: {error}"
                )

        quiz = load_json(
            CURRENT_QUIZ_FILE,
            []
        )

        if quiz:

            st.divider()

            st.write(
                f"**Questions: {len(quiz)}**"
            )

            answers = []

            for index, question in enumerate(
                quiz
            ):

                st.markdown(
                    f"### Question {index + 1}"
                )

                st.write(
                    question.get(
                        "question",
                        ""
                    )
                )

                question_type = question.get(
                    "type",
                    ""
                )

                if question_type == "MCQ":

                    options = question.get(
                        "options",
                        []
                    )

                    answer = st.radio(

                        "Select your answer",

                        options,

                        key=f"mcq_{index}"

                    )

                elif question_type == "True/False":

                    answer = st.radio(

                        "Select your answer",

                        [
                            "True",
                            "False"
                        ],

                        key=f"tf_{index}"

                    )

                else:

                    answer = st.text_input(

                        "Your answer",

                        key=f"answer_{index}"

                    )

                answers.append(
                    answer
                )

            st.divider()

            if st.button(
                "✅ Submit Quiz"
            ):

                if "evaluate" in modules:

                    try:

                        result = modules[
                            "evaluate"
                        ](
                            "",
                            answers
                        )

                        st.session_state.quiz_result = (
                            result
                        )

                    except Exception as error:

                        st.error(
                            f"Evaluation error: {error}"
                        )

                else:

                    st.error(
                        "evaluator.py could not be loaded."
                    )

            result = st.session_state.quiz_result

            if result:

                if result.get(
                    "status"
                ) == "success":

                    score = result.get(
                        "score",
                        0
                    )

                    # Prevent duplicate progress updates
                    # because evaluator.py may already update it.
                    st.success(
                        f"🎉 Score: {score}%"
                    )

                    col1, col2, col3 = st.columns(3)

                    with col1:

                        st.metric(
                            "Correct",
                            result.get(
                                "correct",
                                0
                            )
                        )

                    with col2:

                        st.metric(
                            "Wrong",
                            result.get(
                                "wrong",
                                0
                            )
                        )

                    with col3:

                        st.metric(
                            "Total",
                            result.get(
                                "total",
                                0
                            )
                        )

                    st.info(
                        result.get(
                            "feedback",
                            ""
                        )
                    )

                    with st.expander(
                        "View Question Review"
                    ):

                        for item in result.get(
                            "results",
                            []
                        ):

                            if item.get(
                                "result"
                            ) == "Correct":

                                st.success(
                                    f"Q{item['question_no']}: "
                                    "Correct"
                                )

                            else:

                                st.error(
                                    f"Q{item['question_no']}: "
                                    "Incorrect"
                                )

                            st.write(
                                f"Your answer: "
                                f"{item.get('user_answer', '')}"
                            )

                            st.write(
                                f"Correct answer: "
                                f"{item.get('correct_answer', '')}"
                            )


# =========================================================
# FLASHCARDS
# =========================================================

elif page == "🗂️ Flashcards":

    st.subheader(
        "🗂️ Flashcards"
    )

    text = st.session_state.extracted_text

    if not text:

        st.warning(
            "Please upload study material first."
        )

    elif "flashcards" not in modules:

        st.error(
            "flashcards.py could not be loaded."
        )

    else:

        if st.button(
            "Generate Flashcards"
        ):

            try:

                result = modules[
                    "flashcards"
                ](text)

                st.success(
                    "Flashcards generated."
                )

                st.text(
                    result
                )

            except Exception as error:

                st.error(
                    f"Flashcard error: {error}"
                )


# =========================================================
# REVISION PLANNER
# =========================================================

elif page == "📅 Revision Planner":

    st.subheader(
        "📅 Revision Planner"
    )

    text = st.session_state.extracted_text

    if not text:

        st.warning(
            "Please upload study material first."
        )

    elif "planner" not in modules:

        st.error(
            "planner.py could not be loaded."
        )

    else:

        days = st.slider(

            "How many days do you have?",

            min_value=1,

            max_value=30,

            value=7

        )

        if st.button(
            "Generate Revision Plan"
        ):

            try:

                plan = modules[
                    "planner"
                ](

                    text,

                    days,

                    st.session_state.current_file
                    or "Study Material"

                )

                st.success(
                    "Revision plan generated."
                )

                st.text(
                    plan
                )

            except Exception as error:

                st.error(
                    f"Planner error: {error}"
                )


# =========================================================
# EXPLAIN TOPIC
# =========================================================

elif page == "💡 Explain Topic":

    st.subheader(
        "💡 Explain a Topic"
    )

    text = st.session_state.extracted_text

    if not text:

        st.warning(
            "Please upload study material first."
        )

    elif "explain" not in modules:

        st.error(
            "explain.py could not be loaded."
        )

    else:

        concept = st.text_input(
            "Enter the topic you want explained",
            placeholder="Example: Neural Networks"
        )

        if st.button(
            "💡 Explain"
        ):

            if not concept.strip():

                st.warning(
                    "Please enter a topic."
                )

            else:

                try:

                    explanation = modules[
                        "explain"
                    ](
                        text,
                        concept
                    )

                    st.success(
                        "Explanation found."
                    )

                    st.text(
                        explanation
                    )

                except Exception as error:

                    st.error(
                        f"Explanation error: {error}"
                    )


# =========================================================
# PROGRESS & ANALYTICS
# =========================================================

elif page == "📊 Progress & Analytics":

    st.subheader(
        "📊 Progress & Analytics"
    )

    progress = load_json(
        PROGRESS_FILE,
        DEFAULT_PROGRESS
    )

    col1, col2, col3, col4 = st.columns(4)

    with col1:

        st.metric(
            "Study Days",
            progress.get(
                "study_days",
                0
            )
        )

    with col2:

        st.metric(
            "Files",
            progress.get(
                "total_files_uploaded",
                0
            )
        )

    with col3:

        st.metric(
            "Quizzes",
            progress.get(
                "total_quizzes",
                0
            )
        )

    with col4:

        st.metric(
            "Average Score",
            f"{progress.get('average_score', 0)}%"
        )

    st.divider()

    col1, col2, col3 = st.columns(3)

    with col1:

        st.metric(
            "🏆 Best Score",
            f"{progress.get('best_score', 0)}%"
        )

    with col2:

        st.metric(
            "📝 Last Score",
            f"{progress.get('last_score', 0)}%"
        )

    with col3:

        st.metric(
            "⏱️ Study Time",
            f"{progress.get('study_time_minutes', 0)} min"
        )

    scores = progress.get(
        "quiz_scores",
        []
    )

    if scores:

        st.divider()

        st.subheader(
            "📈 Quiz Score History"
        )

        figure = plt.figure(
            figsize=(9, 4)
        )

        plt.plot(
            range(1, len(scores) + 1),
            scores,
            marker="o"
        )

        plt.title(
            "Quiz Score History"
        )

        plt.xlabel(
            "Quiz Number"
        )

        plt.ylabel(
            "Score (%)"
        )

        plt.ylim(
            0,
            100
        )

        plt.grid(
            True
        )

        st.pyplot(
            figure
        )

        plt.close(
            figure
        )

    else:

        st.info(
            "Complete a quiz to see your score chart."
        )

    st.divider()

    st.subheader(
        "📚 Uploaded Materials"
    )

    files = progress.get(
        "files",
        []
    )

    if files:

        for filename in files:

            st.write(
                f"📄 {filename}"
            )

    else:

        st.info(
            "No files uploaded yet."
        )

    st.divider()

    st.subheader(
        "🔑 Completed Topics"
    )

    topics = progress.get(
        "completed_topics",
        []
    )

    if topics:

        for topic in topics:

            st.write(
                f"✅ {topic}"
            )

    else:

        st.info(
            "No topics completed yet."
        )

    st.divider()

    if st.button(
        "⚠️ Reset Progress"
    ):

        save_json(
            PROGRESS_FILE,
            DEFAULT_PROGRESS
        )

        st.session_state.quiz_result = None

        st.success(
            "Progress has been reset."
        )

        st.rerun()


# =========================================================
# FOOTER
# =========================================================

st.divider()

st.caption(
    "📚 Study Buddy AI Agent | "
    "Python 3.14.7 | "
    "PDF • TXT • CSV • PPT • PPTX"
)