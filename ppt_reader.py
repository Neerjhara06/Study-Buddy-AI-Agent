"""
ppt_reader.py
---------------------------------------
Reads PowerPoint (.pptx) files.

Compatible with:
- Python 3.14.7
- python-pptx
"""

from pathlib import Path
from pptx import Presentation


def read_ppt(file_path: str) -> str:
    """
    Read a PowerPoint presentation and return all text.

    Parameters
    ----------
    file_path : str

    Returns
    -------
    str
    """

    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"{file_path} not found.")

    # python-pptx only supports .pptx
    if path.suffix.lower() == ".ppt":
        raise ValueError(
            "Legacy '.ppt' files are not supported.\n"
            "Please save or convert the presentation to '.pptx' and upload it again."
        )

    try:
        presentation = Presentation(str(path))

    except Exception as e:
        raise RuntimeError(
            f"Unable to open PowerPoint file.\n{e}"
        )

    extracted_text = []

    extracted_text.append("=" * 60)
    extracted_text.append("POWERPOINT PRESENTATION")
    extracted_text.append("=" * 60)

    for slide_number, slide in enumerate(
        presentation.slides,
        start=1
    ):

        extracted_text.append(f"\n\nSLIDE {slide_number}")
        extracted_text.append("-" * 40)

        slide_has_text = False

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:

                    extracted_text.append(text)
                    slide_has_text = True

        if not slide_has_text:
            extracted_text.append("(No text found on this slide.)")

    result = "\n".join(extracted_text).strip()

    if not result:
        raise ValueError(
            "No readable text found in this presentation."
        )

    return result


# ----------------------------------------------------
# Test
# ----------------------------------------------------

if __name__ == "__main__":

    sample = "sample_notes/AI_Notes.pptx"

    try:

        content = read_ppt(sample)

        print(content)

    except Exception as error:

        print(error)